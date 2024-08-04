import requests
from bs4 import BeautifulSoup
import gspread
from pptx import Presentation
from pptx.util import Pt
import tkinter as tk
from tkinter import messagebox
import os
import sys

class Article:
    def __init__(self, num, date, title, link):
        self.num = num
        self.date = date
        self.title = title
        self.link = link

class PostingArticle:
    def __init__(self, row_index, title, link): 
        self.row_index = row_index
        self.title = title
        self.link = link

#정확하게 해당 문자열만 뒤에서 지워주는 함수
def rstrip_exact(s, suffix):
    if suffix and s.endswith(suffix):
        return s[:-len(suffix)]
    return s

class Board:
    def __init__(self, title, link):
        self.title = title
        self.link = link
        self.list = []
        self.new_list = []
        self.new_list_alt = []
        self.top_list = []
        self.to_notify_list = []

    def __str__(self):
        print(f"<{self.title}>")
        for article in self.list[:3]:
            print(f"{article.num} / {article.date} / {article.title} / {article.link}")
        print("")

    def scrap_board(self):
        soup = make_soup(self.link)
        articles = soup.find("tbody").find_all("tr")
        for article in articles:
            contents = article.find_all('td')
            num = contents[0].text.strip()
            title = contents[1].find('a').text
            link = rstrip_exact(self.link + contents[1].find('a')['href'],"article.offset=0&articleLimit=10&totalNoticeYn=N&totalBoardNo=")
            date = contents[-1].text
            article_obj = Article(
                num=num,
                date=date,
                title=title,
                link=link
            )
            self.list.append(article_obj) 

    def update(self, last_num):
        for article in self.list:
            if article.num == last_num:
                index = self.list.index(article)
                self.new_list = self.list[:index]
                break
        for article in self.new_list:
            if article.num == None:
                self.top_list_alt.append([article.num, article.date, article.title, None, None, None, None, article.link])
            else:
                self.new_list_alt.append([article.num, article.date, article.title, None, None, None, None, article.link])
    
    def update_spreadsheet(self):
        worksheet = doc.worksheet(self.title)
        data = worksheet.get_all_values()
        last_num = None
        end_index = 3

        # Find the last_num
        for i in range(2, len(data)):
            if data[i][0].isdigit():
                last_num = data[i][0]
                end_index = i + 1
                break

        # Delete old rows
        worksheet.delete_rows(3, end_index - 1)
        
        # Update and insert new rows
        self.update(last_num)
        worksheet.insert_rows(self.top_list, 3)
        worksheet.insert_rows(self.new_list_alt, 3)

    def update_to_notify_list(self):
        self.to_notify_list = []
        worksheet = doc.worksheet(self.title)
        data = worksheet.get_all_values()
        for i in range(2, len(data)):
            if data[i][6] == 'FALSE':
                self.to_notify_list.append(PostingArticle(i, data[i][2], data[i][7]))
            elif data[i][6] == 'TRUE':
                break
    
    def write(self, file):
        for article in self.to_notify_list:
            file.write(f"📌 {article.title}\n")
            file.write(f"🔗 링크\n{article.link}\n\n")
            
def load_spreadsheet():
    if getattr(sys, 'frozen', False):  # Check if running as a compiled executable
        json_file_path = os.path.join(os.path.dirname(sys.executable), "geuruteogi-automation-f1dddf13478a.json")
    else:
        json_file_path = os.path.join(os.path.dirname(__file__), "geuruteogi-automation-f1dddf13478a.json")

    gc = gspread.service_account(json_file_path)
    spreadsheet_url = url_entry.get()
    global doc
    doc = gc.open_by_url(spreadsheet_url)


def fetch_html(url):
    response = requests.get(url)
    if response.status_code == 200:
        return response.text
    else:
        raise Exception(f"Failed to fetch the URL: {url}, Status code: {response.status_code}")

def make_soup(url):
    html = fetch_html(url)
    return BeautifulSoup(html, 'html.parser')

def update(boards):
    for board in boards.values():
        board.scrap_board()
        board.update_spreadsheet()

def write_announcement(boards, month, week):
    for board in boards.values():
        board.update_to_notify_list()

    with open('announcement.txt', 'w', encoding='utf-8') as file:
        file.write(f"[2024 {month}월 {week}주차 전공소식 공유]\n") 
        file.write("\n안녕하세요.\n정보대학 제8대 학생회 'Canvas' 입니다.\n")
        file.write("금주의 전공소식을 공유드리니 많은 관심 부탁드립니다. 공지글의 링크로 들어가시면 더 자세한 정보를 확인하실 수 있습니다.\n")
        
        file.write("\n📢 일반공지\n")
        boards['공지사항'].write(file)
        boards['장학공지'].write(file)

        file.write("\n📍 행사 및 공모전\n")
        boards['진로정보(공모전)'].write(file)
        boards['행사 및 소식'].write(file)

        file.write("\n💼 채용 및 인턴 모집\n")
        boards['진로정보(채용)'].write(file)
        boards['진로정보(인턴)'].write(file)

        file.write("\n📘 전공 관련 교육행사\n")
        boards['진로정보(교육)'].write(file)

def select_shape_by_text(slide, text):
    shapes = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                if text in paragraph.text:
                    shapes.append(shape)
    return shapes

def copy_slide(pres, layout_name):
    slide_layout = None
    for layout in pres.slide_layouts:
        if layout.name == layout_name:
            slide_layout = layout
            break
    if slide_layout is None:
        raise ValueError(f"Layout {layout_name} not found in presentation.")
    slide = pres.slides.add_slide(slide_layout)
    return slide

def list_chunk(lst, n=5):
    return [lst[i:i+n] for i in range(0, len(lst), n)]

def make_ppt(boards, month, week, name, template_path='template.pptx'):
    prs = Presentation(template_path)

    for board in boards.values():
        board.update_to_notify_list()

    title_page = prs.slides[0]
    date_shape = select_shape_by_text(title_page, '2024년 n월 m주차')[0]
    date_shape.text = f"2024년 {month}월 {week}주차"
    author_shape = select_shape_by_text(title_page, "담당자 : 교육/진로국원 OOO")[0]
    author_shape.text_frame.paragraphs[0].runs[0].text = f"담당자 : 교육/진로국원 {name}"

    list = boards['공지사항'].to_notify_list + boards['장학공지'].to_notify_list + boards['행사 및 소식'].to_notify_list + boards['진로정보(공모전)'].to_notify_list + boards['진로정보(채용)'].to_notify_list + boards['진로정보(인턴)'].to_notify_list + boards['진로정보(교육)'].to_notify_list
    list = list_chunk(list)

    slide_index = 1
    for sub_list in list:
        article_index = 0
        shapes = select_shape_by_text(prs.slides[slide_index], "제목제목제목제목제목제목제목제목")
        for article in sub_list:
            print(f"slide: {slide_index}, article: {article_index}")
            shape = shapes[article_index]
            shape.text_frame.paragraphs[0].runs[0].text = article.title
            article_index += 1
        slide_index += 1

    prs.save(f'전공소식공유_{month}월_{week}주차.pptx')

boards = {
    '공지사항': Board('공지사항', 'https://info.korea.ac.kr/info/board/notice_under.do'),
    '장학공지': Board('장학공지', 'https://info.korea.ac.kr/info/board/scholarship_under.do'),
    '행사 및 소식': Board('행사 및 소식', 'https://info.korea.ac.kr/info/board/news.do'),
    '진로정보(채용)': Board("진로정보(채용)", 'https://info.korea.ac.kr/info/board/course_job.do'),
    '진로정보(교육)': Board("진로정보(교육)", 'https://info.korea.ac.kr/info/board/course_program.do'),
    '진로정보(인턴)': Board("진로정보(인턴)", 'https://info.korea.ac.kr/info/board/course_intern.do'),
    '진로정보(공모전)': Board("진로정보(공모전)", 'https://info.korea.ac.kr/info/board/course_competition.do'),
}

# GUI Implementation

def update_action():
    load_spreadsheet()
    update(boards)
    messagebox.showinfo("Update", "Boards updated successfully")

def write_action():
    load_spreadsheet()
    month = month_entry.get()
    week = week_entry.get()
    write_announcement(boards, month, week)
    messagebox.showinfo("Write", f"Announcement written for {month}월 {week}주차")

def ppt_action():
    load_spreadsheet()
    name = name_entry.get()
    month = month_entry.get()
    week = week_entry.get()
    make_ppt(boards, month, week, name)
    messagebox.showinfo("PPT", f"PPT created for {month}월 {week}주차 by {name}")

app = tk.Tk()
app.title("KUCI Automation")

tk.Label(app, text="Spreadsheet URL:").grid(row=0, column=0, padx=10, pady=10)
url_entry = tk.Entry(app)
url_entry.grid(row=0, column=1, padx=10, pady=10)
url_entry.insert(0, "https://docs.google.com/spreadsheets/d/1KQ5nffRAkpYr22HE5P43WiYU8sCySBW8I1IVP77Elpk/edit?usp=sharing")

tk.Label(app, text="Month:").grid(row=1, column=0, padx=10, pady=10)
month_entry = tk.Entry(app)
month_entry.grid(row=1, column=1, padx=10, pady=10)

tk.Label(app, text="Week:").grid(row=2, column=0, padx=10, pady=10)
week_entry = tk.Entry(app)
week_entry.grid(row=2, column=1, padx=10, pady=10)

tk.Label(app, text="Name:").grid(row=3, column=0, padx=10, pady=10)
name_entry = tk.Entry(app)
name_entry.grid(row=3, column=1, padx=10, pady=10)

tk.Button(app, text="Update", command=update_action).grid(row=4, column=0, padx=10, pady=10)
tk.Button(app, text="Write", command=write_action).grid(row=4, column=1, padx=10, pady=10)
tk.Button(app, text="PPT", command=ppt_action).grid(row=4, column=2, padx=10, pady=10)
tk.Button(app, text="Exit", command=app.quit).grid(row=5, column=1, padx=10, pady=10)

app.mainloop()
