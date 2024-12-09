from pptx import Presentation
from models import Board
import os
import itertools

def chunk_list(my_list, chunk_size=4):
    return [my_list[i:i+chunk_size] for i in range(0, len(my_list), chunk_size)]

def select_shape_by_text(slide, text):
    shapes = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                if text in paragraph.text:
                    shapes.append(shape)
    return shapes

def write_slide(slide, articles):
    shapes = list(itertools.islice(slide.shapes, 1, 8, 2)) # 제목들만
    for shape, article in zip(shapes, articles):
        shape.text_frame.word_wrap = True # 줄바꿈 허용
        shape.text_frame.paragraphs[0].runs[0].text = article.title

        # 줄바꿈 되지 않을 때까지 font size를 줄이는 코드를 짜려 했지만 현재로선 실패
        # text_frame = shape.text_frame
        # text_frame.word_wrap = True # 줄바꿈 허용
        # paragraph = text_frame.paragraphs[0]

        # run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
        # run.text = article.title

        # if len(text_frame.paragraphs) > 1:
        #     for paragraph in text_frame.paragraphs:
        #         paragraph.runs[0].font.size -= Pt(12)
    return

# 글자크기, 글꼴 유지하면서 슬라이드 복제하기
def duplicate_slide(prs, slide):
    # Create a new slide with the same layout as the original slide
    new_slide = prs.slides.add_slide(slide.slide_layout)

    # Copy all shapes from the original slide to the new slide
    for shape in slide.shapes:
        if shape.is_placeholder:
            # Copy placeholder text and font properties
            new_shape = new_slide.placeholders[shape.placeholder_format.idx]
            if shape.has_text_frame:
                new_shape.text_frame.text = shape.text_frame.text
                for p, new_p in zip(shape.text_frame.paragraphs, new_shape.text_frame.paragraphs):
                    for r, new_r in zip(p.runs, new_p.runs):
                        if r.font.size:
                            new_r.font.size = r.font.size
                        if r.font.name:
                            new_r.font.name = r.font.name
        elif shape.has_text_frame:
            # Copy regular text boxes and font properties
            textbox = new_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
            textbox.text = shape.text_frame.text
            new_text_frame = textbox.text_frame
            for p, new_p in zip(shape.text_frame.paragraphs, new_text_frame.paragraphs):
                for r, new_r in zip(p.runs, new_p.runs):
                    if r.font.size:
                        new_r.font.size = r.font.size
                    if r.font.name:
                        new_r.font.name = r.font.name

    # placeholder 삭제하기 (있다면)
    if new_slide.shapes and new_slide.shapes[0].is_placeholder:
        new_slide.shapes[0]._element.getparent().remove(new_slide.shapes[0]._element)

    return new_slide   

def make_ppt(month, week, writer, template_path='template.pptx'):
    template_path = os.path.join(os.path.dirname(__file__), template_path)
    print(template_path)
    prs = Presentation(template_path)

    title_page = prs.slides[0]
    date_shape = select_shape_by_text(title_page, 'n월 m주차')[0]
    date_shape.text = f"2024년 {month}월 {week}주차"
    author_shape = select_shape_by_text(title_page, "담당자")[0]
    author_shape.text_frame.paragraphs[0].runs[0].text = f"담당자: 교육진로국원 {writer}"

    template_page = prs.slides[1]

    ppt_sectors = {
        '일반공지': Board.query.filter_by(name='공지사항').first().articles,
        '행사 및 공모전': Board.query.filter_by(name='행사 및 소식').first().articles + Board.query.filter_by(name='진로정보(공모전)').first().articles,
        '채용 및 인턴 모집': Board.query.filter_by(name='진로정보(채용)').first().articles + Board.query.filter_by(name='진로정보(인턴)').first().articles,
        '전공 관련 교육행사': Board.query.filter_by(name='진로정보(교육)').first().articles
    }

    for sector_title, articles in ppt_sectors.items():
        article_chunks = chunk_list(articles)
 
        for article_chunk in article_chunks:
            slide = duplicate_slide(prs, template_page)
            slide.shapes[0].text_frame.paragraphs[0].runs[0].text = sector_title
            write_slide(slide, article_chunk)

    # Save the modified presentation
    file_path = f"{month}월_{week}주차_전공소식공유.pptx"
    prs.save(file_path)
    
    return file_path